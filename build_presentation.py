#!/usr/bin/env python3
"""
build_presentation.py

Generates a PowerPoint presentation showing simulation results, methodology,
calibrated parameters, and state-level comparative metrics for all 153 counties.
"""

import os
import sys
import glob
import yaml
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Hardcoded list of all 153 counties from run_all_sims.sh
COUNTIES = [
    '01009', '01031', '01039', '01045', '01049', '01071', '13013', '13071', '13103', '13115',
    '13137', '13153', '17011', '17021', '17037', '17055', '17063', '17073', '17091', '17095',
    '17103', '17115', '17117', '17121', '17141', '17167', '19017', '19027', '19113', '19155',
    '19169', '20035', '20045', '20061', '20079', '20103', '20125', '20149', '20161', '20169',
    '21009', '21029', '21035', '21071', '21073', '21083', '21089', '21093', '21113', '21179',
    '21199', '21209', '21211', '22001', '22005', '22039', '22063', '26023', '26027', '26041',
    '26057', '26059', '26067', '26073', '26123', '27005', '27027', '27041', '27059', '27085',
    '27109', '29021', '29027', '29051', '29101', '30029', '30047', '30049', '30093', '36003',
    '36011', '36013', '36031', '36051', '38059', '38077', '38093', '39005', '39007', '39011',
    '39013', '39015', '39021', '39023', '39031', '39033', '39037', '39039', '39045', '39051',
    '39057', '39059', '39063', '39071', '39077', '39079', '39083', '39087', '40021', '40089',
    '40097', '40121', '40147', '42005', '42013', '42025', '42031', '42051', '42055', '42059',
    '42061', '42075', '42085', '42087', '46011', '46013', '46029', '46035', '46081', '46083',
    '46135', '47003', '47005', '47009', '47017', '48001', '48005', '48013', '48041', '48049',
    '48071', '48091', '48097', '48143', '48147', '48181', '48189', '48213', '51047', '51095',
    '51137', '51149', '51161'
]

# Standard state FIPS dictionary mapping
STATE_FIPS_MAP = {
    '01': 'AL', '02': 'AK', '04': 'AZ', '05': 'AR', '06': 'CA', '08': 'CO', '09': 'CT', '10': 'DE',
    '11': 'DC', '12': 'FL', '13': 'GA', '15': 'HI', '16': 'ID', '17': 'IL', '18': 'IN', '19': 'IA',
    '20': 'KS', '21': 'KY', '22': 'LA', '23': 'ME', '24': 'MD', '25': 'MA', '26': 'MI', '27': 'MN',
    '28': 'MS', '29': 'MO', '30': 'MT', '31': 'NE', '32': 'NV', '33': 'NH', '34': 'NJ', '35': 'NM',
    '36': 'NY', '37': 'NC', '38': 'ND', '39': 'OH', '40': 'OK', '41': 'OR', '42': 'PA', '44': 'RI',
    '45': 'SC', '46': 'SD', '47': 'TN', '48': 'TX', '49': 'UT', '50': 'VT', '51': 'VA', '53': 'WA',
    '54': 'WV', '55': 'WI', '56': 'WY'
}

def load_county_names():
    """Builds a mapping from 5-digit FIPS to (County Name, State Abbreviation)."""
    fips_map = {}
    policy_db_path = os.path.join("data", "county_policy_data", "39109-0001-Data.tsv")
    if os.path.exists(policy_db_path):
        try:
            df = pd.read_csv(policy_db_path, sep="\t", usecols=["COUNTY_FIPS", "COUNTYNAME", "STATE"])
            df = df.dropna(subset=["COUNTY_FIPS"]).drop_duplicates(subset=["COUNTY_FIPS"])
            for _, row in df.iterrows():
                fips_str = str(int(row["COUNTY_FIPS"])).zfill(5)
                fips_map[fips_str] = (str(row["COUNTYNAME"]).strip(), str(row["STATE"]).strip())
        except Exception as e:
            print(f"Warning loading policy dataset for county names: {e}")

    return fips_map

def load_config():
    """Loads configuration metadata from covid_abm/yamls/config.yaml or covid_abm/config.yaml."""
    config_path = os.path.join("covid_abm", "yamls", "config.yaml")
    if not os.path.exists(config_path):
        config_path = os.path.join("covid_abm", "config.yaml")
    
    if not os.path.exists(config_path):
        raise FileNotFoundError("Could not locate config.yaml in covid_abm/yamls/ or covid_abm/.")
        
    with open(config_path, "r") as f:
        return yaml.safe_load(f)

def get_county_folder_and_graph(population, config):
    """
    Constructs the config directory for a county and finds the latest epoch's simulation_results.png
    and calibrated_params.txt.
    """
    meta = config.get("simulation_metadata", {})
    date = meta.get("DATE", "202010-202104")
    initial_rate = meta.get("INITIAL_INFECTION_RATE", 0.0005)
    exposed_to_infected = meta.get("EXPOSED_TO_INFECTED_TIME", 3)
    infected_to_recovered = meta.get("INFECTED_TO_RECOVERED_TIME", 5)
    with_k = meta.get("WITH_K", True)
    with_vacc = meta.get("WITH_VACC", False)
    use_7day_avg = meta.get("USE_7DAY_AVG", True)
    metro_phase = meta.get("metro_calibration_phase", 0)

    # Reconstruct folder string based on abm_nets.py
    folder_pattern = f"{initial_rate}_{exposed_to_infected}_{infected_to_recovered}_{with_k}_{with_vacc}_{use_7day_avg}_metro_{metro_phase}"
    base_dir = os.path.join("result_graphs", population, date, folder_pattern)

    # Check if exact directory exists, or search for potential metro phase variations
    if not os.path.exists(base_dir):
        matches = glob.glob(os.path.join("result_graphs", population, date, f"{initial_rate}_{exposed_to_infected}_{infected_to_recovered}_*"))
        if matches:
            base_dir = sorted(matches)[-1]

    if not os.path.exists(base_dir):
        return None, None, None

    calibrated_params_file = os.path.join(base_dir, "calibrated_params.txt")

    subdirs = [d for d in os.listdir(base_dir) if os.path.isdir(os.path.join(base_dir, d)) and d.isdigit()]
    subdirs = sorted(subdirs, key=lambda x: int(x), reverse=True)

    graph_path = None
    latest_epoch = None

    for epoch_dir in subdirs:
        img_path = os.path.join(base_dir, epoch_dir, "simulation_results.png")
        if os.path.exists(img_path):
            graph_path = img_path
            latest_epoch = epoch_dir
            break

    return base_dir, graph_path, calibrated_params_file

def parse_calibrated_params(params_file):
    """
    Reads calibrated_params.txt and extracts Average R (Scaling Factor), Initial Infection Rate (per 100k), and K parameter.
    """
    if not params_file or not os.path.exists(params_file):
        return None

    try:
        data = np.loadtxt(params_file)
        if data.ndim == 0 or len(data) < 3:
            return None

        weekly_r2 = data[:-2]
        avg_r2 = float(np.mean(weekly_r2))
        initial_rate = float(data[-2])
        k_param = float(data[-1])

        return {
            "weekly_r2": weekly_r2,
            "avg_r2": avg_r2,
            "initial_rate": initial_rate,
            "initial_rate_per_100k": initial_rate * 100000.0,
            "k_param": k_param,
            "num_weeks": len(weekly_r2)
        }
    except Exception as e:
        print(f"Error parsing calibrated params from {params_file}: {e}")
        return None

def add_parameter_explanation_slide(prs):
    """Adds an initial slide explaining each parameter and its formula."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title Header
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Simulation Parameter Definitions & Equations"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 38, 59)

    params_data = [
        {
            "name": "Scaling Factor (R)",
            "desc": "Scales the overall transmission rate λ across agent contact networks. Under simplifying assumptions, it represents the mean number of secondary infections caused by an infectious individual.",
            "formula": "λ(t, s_i, a_s, n) = [ R · S_{a_s} · A_{s_i} · B_n / I_bar ] · ∫_{t-1}^t f_Γ(u; μ_i, σ_i²) du"
        },
        {
            "name": "Initial Infection Rate (per 100k)",
            "desc": "Proportion of individuals initialized in the infected compartment at timestep t=0, expressed per 100,000 residents in the county population (I₀ × 10⁵).",
            "formula": "Initial Infected Agents = ⌈ (Initial Rate per 100k / 100,000) · N ⌉"
        },
        {
            "name": "K Parameter (Underreporting Factor)",
            "desc": "Scalar multiplier adjusting raw simulated new infections (N_t) to observable reported case counts (Ŷ_t) to account for public health surveillance underreporting.",
            "formula": "Ŷ_t = k · N_t"
        }
    ]

    for idx, pinfo in enumerate(params_data):
        top_pos = Inches(1.4 + idx * 1.8)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top_pos, Inches(12.1), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 247, 250)
        card.line.color.rgb = RGBColor(210, 215, 225)
        card.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(Inches(0.8), top_pos + Inches(0.1), Inches(11.7), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p_name = tf.paragraphs[0]
        p_name.text = pinfo["name"]
        p_name.font.size = Pt(18)
        p_name.font.bold = True
        p_name.font.color.rgb = RGBColor(13, 27, 42)

        p_desc = tf.add_paragraph()
        p_desc.text = pinfo["desc"]
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = RGBColor(65, 90, 119)
        p_desc.space_before = Pt(2)

        p_form = tf.add_paragraph()
        p_form.text = f"Formula: {pinfo['formula']}"
        p_form.font.size = Pt(13)
        p_form.font.bold = True
        p_form.font.color.rgb = RGBColor(27, 38, 59)
        p_form.space_before = Pt(3)

def add_methodology_slide(prs):
    """Adds a slide detailing methodology, county selection, and policy data sources."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title Header
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Methodology, Assumptions & Policy Data Sources"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 38, 59)

    sections = [
        {
            "title": "1. County Selection & Timeline",
            "bullets": [
                "Evaluated 153 U.S. counties with population sizes between 20,000 and 200,000.",
                "Simulated over a 168-day ground-truth window (October 26, 2020 – April 11, 2021).",
                "Qualitatively filtered to retain counties where gradient-based optimization converged to an accurate epidemiological baseline."
            ]
        },
        {
            "title": "2. School & Workplace Closures (ICPSR 39109 Policy Dataset)",
            "bullets": [
                "Source: County-Level Policy Database (ICPSR 39109, provided by John; data/county_policy_data/39109-0001-Data.tsv).",
                "Extracted C1_SCHOOL and C2_WORKPLACE mandates (with state-level fallbacks S_C1_SCHOOL / S_C2_WORKPLACE if county values were missing/99).",
                "Binarization: Mandate severity levels 0 & 1 -> 0 (weak/open); levels 2 & 3 -> 1 (strong/closed).",
                "ABM Implementation: Enforced via a Network Freezing Mechanism that probabilistically prunes contact edges (0-25% edge removal for open vs 75-100% for closed)."
            ]
        },
        {
            "title": "3. Vaccination Rollout (CDC Socrata API)",
            "bullets": [
                "Source: CDC COVID-19 Vaccinations County Dataset via Socrata API (https://data.cdc.gov/resource/8xkx-amqh.csv, querying administered_dose1_recip).",
                "Extraction: Ingests daily first-dose recipient increments (new_vax = diff(administered_dose1_recip)).",
                "ABM Implementation: In apply_vaccines, N_vax susceptible agents are probabilistically transitioned directly to the Recovered state (R) each timestep."
            ]
        }
    ]

    for idx, sinfo in enumerate(sections):
        top_pos = Inches(1.3 + idx * 1.9)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top_pos, Inches(12.1), Inches(1.75))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(248, 249, 250)
        box.line.color.rgb = RGBColor(220, 224, 230)
        box.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(Inches(0.8), top_pos + Inches(0.08), Inches(11.7), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = sinfo["title"]
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = RGBColor(13, 27, 42)

        for b in sinfo["bullets"]:
            p_b = tf.add_paragraph()
            p_b.text = f"•  {b}"
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = RGBColor(65, 90, 119)
            p_b.space_before = Pt(2)

def add_state_scaling_factor_slides(prs, county_params_map):
    """
    Generates an individual scaling factor (R_t) trend plot and dedicated slide for each state.
    """
    if not county_params_map:
        print("No county parameters available for state scaling factor plots.")
        return

    # Group weekly R by state
    state_weekly_r = {}
    for fips, info in county_params_map.items():
        state = info["state"]
        weekly_r = info["params"]["weekly_r2"]
        if state not in state_weekly_r:
            state_weekly_r[state] = []
        state_weekly_r[state].append(weekly_r)

    if not state_weekly_r:
        return

    blank_layout = prs.slide_layouts[6]

    for state, r_list in sorted(state_weekly_r.items()):
        plt.figure(figsize=(10, 5.2))
        num_weeks = max(len(w) for w in r_list)
        weeks = np.arange(1, num_weeks + 1)

        # Plot individual county trajectories lightly
        for w in r_list:
            plt.plot(weeks[:len(w)], w, color='#64748B', alpha=0.35, linewidth=1.2)

        # Plot mean state trajectory
        min_len = min(len(w) for w in r_list)
        truncated_r = [w[:min_len] for w in r_list]
        avg_r_state = np.mean(truncated_r, axis=0)
        plt.plot(weeks[:min_len], avg_r_state, color='#1B263B', marker='o', linewidth=2.5, label=f"State Mean ({state})")

        plt.xlabel("Simulation Week", fontsize=12)
        plt.ylabel("Scaling Factor (R)", fontsize=12)
        plt.title(f"Scaling Factor (R) Trend Over Time – {state} ({len(r_list)} counties)", fontsize=14, fontweight='bold')
        plt.grid(True, linestyle='--', alpha=0.6)
        plt.legend(loc='upper right', fontsize=10)
        plt.tight_layout()

        img_path = f"/tmp/state_scaling_factor_{state}.png"
        plt.savefig(img_path, dpi=300)
        plt.close()

        slide = prs.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Scaling Factor (R) Trend Over Time – {state}"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = RGBColor(27, 38, 59)

        slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.3), width=Inches(12.1))

def calculate_county_nrmse(fips, config):
    """Calculates Normalized RMSE (NRMSE) between simulation factual cases and ground truth actual cases."""
    base_dir, graph_path, params_file = get_county_folder_and_graph(fips, config)
    meta = config.get("simulation_metadata", {})
    date = meta.get("DATE", "202010-202104")
    use_7day_avg = meta.get("USE_7DAY_AVG", True)
    initial_rate = meta.get("INITIAL_INFECTION_RATE", 0.0005)
    exposed_to_infected = meta.get("EXPOSED_TO_INFECTED_TIME", 3)
    infected_to_recovered = meta.get("INFECTED_TO_RECOVERED_TIME", 5)

    # Load actual daily data
    daily_csv = os.path.join("data", "processed_data", fips, date, "daily_data.csv")
    if not os.path.exists(daily_csv):
        daily_csv = os.path.join("data", "processed_data", fips, "daily_data.csv")

    if not os.path.exists(daily_csv):
        return None

    try:
        actual_df = pd.read_csv(daily_csv)
        case_col = 'cases_7day_avg' if use_7day_avg and 'cases_7day_avg' in actual_df.columns else 'cases'
        actual_cases = actual_df[case_col].values
    except Exception:
        return None

    # Load simulation factual cases
    sim_cases = None
    gen_fac_path = os.path.join("results", fips, f"{initial_rate}_{exposed_to_infected}_{infected_to_recovered}_metro_0", "generated_factual.csv")
    if os.path.exists(gen_fac_path):
        try:
            gen_df = pd.read_csv(gen_fac_path)
            sim_cases = gen_df["generated_factual_cases"].values
        except Exception:
            pass

    if sim_cases is None and base_dir and os.path.exists(base_dir):
        cf_path = os.path.join(base_dir, "counterfactual_data0.csv")
        if os.path.exists(cf_path):
            try:
                cf_df = pd.read_csv(cf_path)
                if "counterfactual_cases" in cf_df.columns:
                    sim_cases = cf_df["counterfactual_cases"].values
            except Exception:
                pass

    if sim_cases is None or len(actual_cases) == 0:
        return None

    min_len = min(len(sim_cases), len(actual_cases))
    sim_cases = sim_cases[:min_len]
    actual_cases = actual_cases[:min_len]

    mean_actual = np.mean(actual_cases)
    if mean_actual == 0:
        return None

    rmse = np.sqrt(np.mean((sim_cases - actual_cases) ** 2))
    nrmse = rmse / mean_actual
    return nrmse

def add_state_nrmse_slide(prs, config, county_name_lookup):
    """Calculates NRMSE per county, aggregates by state, plots comparison, and adds slide."""
    state_nrmse = {}
    for fips in COUNTIES:
        nrmse = calculate_county_nrmse(fips, config)
        if nrmse is not None:
            state_abbr = county_name_lookup.get(fips, (None, STATE_FIPS_MAP.get(fips[:2], "US")))[1]
            if state_abbr not in state_nrmse:
                state_nrmse[state_abbr] = []
            state_nrmse[state_abbr].append(nrmse)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "State Calibration Performance by Normalized RMSE (NRMSE)"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 38, 59)

    if not state_nrmse:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
        shape.line.color.rgb = RGBColor(210, 215, 225)
        tf_g = shape.text_frame
        p_g = tf_g.paragraphs[0]
        p_g.text = "State NRMSE Comparison Plot\n(Will be generated dynamically when executed on the server with simulation results)"
        p_g.alignment = PP_ALIGN.CENTER
        p_g.font.size = Pt(18)
        p_g.font.color.rgb = RGBColor(120, 120, 120)
        return

    states = sorted(state_nrmse.keys())
    mean_nrmse = [np.mean(state_nrmse[st]) for st in states]
    std_nrmse = [np.std(state_nrmse[st]) if len(state_nrmse[st]) > 1 else 0 for st in states]

    plt.figure(figsize=(10, 5.5))
    bars = plt.bar(states, mean_nrmse, yerr=std_nrmse, capsize=4, color='#1B263B', edgecolor='#0D1B2A', alpha=0.85)
    plt.xlabel("State", fontsize=12)
    plt.ylabel("Normalized RMSE (RMSE / Mean Actual Cases)", fontsize=12)
    plt.title("Calibration Accuracy by State (Lower NRMSE = Better Fit)", fontsize=14, fontweight='bold')
    plt.grid(axis='y', linestyle='--', alpha=0.6)

    for bar in bars:
        height = bar.get_height()
        plt.text(bar.get_x() + bar.get_width()/2., height + 0.01, f"{height:.3f}", ha='center', va='bottom', fontsize=9)

    plt.tight_layout()
    img_path = "/tmp/state_nrmse_comparison.png"
    plt.savefig(img_path, dpi=300)
    plt.close()

    slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.3), width=Inches(12.1))

def create_presentation(output_pptx="simulation_results_presentation.pptx"):
    """Builds the full presentation for all counties."""
    config = load_config()
    county_name_lookup = load_county_names()

    prs = Presentation()
    # Set 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    print("Adding introductory slides (Parameter Definitions & Methodology)...")
    add_parameter_explanation_slide(prs)
    add_methodology_slide(prs)

    county_params_map = {}

    print(f"Generating presentation for {len(COUNTIES)} counties...")

    for idx, fips in enumerate(COUNTIES, 1):
        if fips in county_name_lookup:
            c_name, state_abbr = county_name_lookup[fips]
            title_text = f"{c_name}, {state_abbr} (FIPS: {fips})"
        else:
            state_abbr = STATE_FIPS_MAP.get(fips[:2], "US")
            c_name = f"County {fips}"
            title_text = f"{c_name}, {state_abbr} (FIPS: {fips})"

        slide = prs.slides.add_slide(blank_layout)

        # -------------------------------------------------------------
        # Slide Header (Top Left)
        # -------------------------------------------------------------
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(27, 38, 59)

        # -------------------------------------------------------------
        # Left Side: Fitting Graph (simulation_results.png)
        # -------------------------------------------------------------
        base_dir, graph_path, params_file = get_county_folder_and_graph(fips, config)

        if graph_path and os.path.exists(graph_path):
            slide.shapes.add_picture(graph_path, Inches(0.6), Inches(1.4), width=Inches(7.2))
        else:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4), Inches(7.2), Inches(5.4))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            shape.line.color.rgb = RGBColor(200, 200, 200)
            tf_g = shape.text_frame
            p_g = tf_g.paragraphs[0]
            p_g.text = f"Fitting Graph (simulation_results.png)\nNot Found for {fips}"
            p_g.alignment = PP_ALIGN.CENTER
            p_g.font.size = Pt(18)
            p_g.font.color.rgb = RGBColor(120, 120, 120)

        # -------------------------------------------------------------
        # Right Side: Calibrated Parameters Card
        # -------------------------------------------------------------
        card_left = Inches(8.2)
        card_top = Inches(1.4)
        card_width = Inches(4.5)
        card_height = Inches(5.4)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 249, 250)
        card.line.color.rgb = RGBColor(220, 224, 230)
        card.line.width = Pt(1.5)

        card_title_box = slide.shapes.add_textbox(card_left + Inches(0.2), card_top + Inches(0.2), card_width - Inches(0.4), Inches(0.6))
        tf_ct = card_title_box.text_frame
        tf_ct.word_wrap = True
        p_ct = tf_ct.paragraphs[0]
        p_ct.text = "Final Calibrated Parameters"
        p_ct.font.size = Pt(20)
        p_ct.font.bold = True
        p_ct.font.color.rgb = RGBColor(13, 27, 42)

        params = parse_calibrated_params(params_file)

        param_box = slide.shapes.add_textbox(card_left + Inches(0.2), card_top + Inches(0.85), card_width - Inches(0.4), card_height - Inches(1.0))
        tf_p = param_box.text_frame
        tf_p.word_wrap = True

        if params:
            county_params_map[fips] = {
                "county_name": c_name,
                "state": state_abbr,
                "params": params
            }

            # 1. Scaling Factor (R)
            p_label1 = tf_p.paragraphs[0]
            p_label1.text = "Scaling Factor (R)"
            p_label1.font.size = Pt(14)
            p_label1.font.bold = True
            p_label1.font.color.rgb = RGBColor(65, 90, 119)

            p_val1 = tf_p.add_paragraph()
            p_val1.text = f"{params['avg_r2']:.4f}"
            p_val1.font.size = Pt(20)
            p_val1.font.bold = True
            p_val1.font.color.rgb = RGBColor(27, 38, 59)
            p_val1.space_after = Pt(12)

            # 2. Initial Infection Rate (per 100k)
            p_label2 = tf_p.add_paragraph()
            p_label2.text = "Initial Infection Rate (per 100k)"
            p_label2.font.size = Pt(14)
            p_label2.font.bold = True
            p_label2.font.color.rgb = RGBColor(65, 90, 119)

            p_val2 = tf_p.add_paragraph()
            p_val2.text = f"{params['initial_rate_per_100k']:.2f} per 100k"
            p_val2.font.size = Pt(20)
            p_val2.font.bold = True
            p_val2.font.color.rgb = RGBColor(27, 38, 59)
            p_val2.space_after = Pt(12)

            # 3. K Parameter (Underreporting Factor)
            p_label3 = tf_p.add_paragraph()
            p_label3.text = "K Parameter (Underreporting)"
            p_label3.font.size = Pt(14)
            p_label3.font.bold = True
            p_label3.font.color.rgb = RGBColor(65, 90, 119)

            p_val3 = tf_p.add_paragraph()
            p_val3.text = f"{params['k_param']:.4f}"
            p_val3.font.size = Pt(20)
            p_val3.font.bold = True
            p_val3.font.color.rgb = RGBColor(27, 38, 59)

        else:
            p_err = tf_p.paragraphs[0]
            p_err.text = "Calibrated parameters (calibrated_params.txt)\nnot found or incomplete."
            p_err.font.size = Pt(14)
            p_err.font.italic = True
            p_err.font.color.rgb = RGBColor(180, 50, 50)

        if idx % 20 == 0 or idx == len(COUNTIES):
            print(f"Processed {idx}/{len(COUNTIES)} county slides...")

    print("Adding state-level individual scaling factor trend slides...")
    add_state_scaling_factor_slides(prs, county_params_map)
    
    print("Adding state NRMSE comparison slide...")
    add_state_nrmse_slide(prs, config, county_name_lookup)

    prs.save(output_pptx)
    print(f"\nSuccessfully generated presentation: {os.path.abspath(output_pptx)}")

if __name__ == "__main__":
    out_file = sys.argv[1] if len(sys.argv) > 1 else "simulation_results_presentation.pptx"
    create_presentation(out_file)
